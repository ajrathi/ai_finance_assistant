"""
Generate the Finnie AI Finance Assistant team demo presentation.
Run: /tmp/finnie_test_env/bin/python create_presentation.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand Colors ──────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x1B, 0x3E)   # dark background
TEAL       = RGBColor(0x00, 0xB4, 0xD8)   # accent / highlight
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF4, 0xF8)
DARK_GRAY  = RGBColor(0x44, 0x55, 0x66)
GREEN      = RGBColor(0x2E, 0xCC, 0x71)
ORANGE     = RGBColor(0xF3, 0x96, 0x16)
SLATE      = RGBColor(0x1E, 0x3A, 0x5F)

W, H = Inches(13.33), Inches(7.5)   # 16:9 widescreen

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H


# ── Helpers ───────────────────────────────────────────────────────────────────

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width or 1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h, font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_multiline(slide, lines, l, t, w, h, font_size=14, color=WHITE,
                  bold=False, align=PP_ALIGN.LEFT, line_spacing=None):
    """Add a text box with multiple paragraphs."""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txb


def slide_header(slide, title, subtitle=None, title_color=WHITE, sub_color=TEAL):
    add_rect(slide, 0, 0, 13.33, 1.2, SLATE)
    add_text(slide, title, 0.4, 0.15, 12, 0.7, font_size=28, bold=True,
             color=title_color, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.75, 12, 0.4, font_size=14,
                 color=sub_color, align=PP_ALIGN.LEFT)


def bullet_box(slide, items, l, t, w, h, font_size=15, bullet="▸", color=WHITE, spacing=22):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = Pt(spacing)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txb


def footer(slide, text="Finnie — AI Finance Assistant  |  Educational AI, not financial advice"):
    add_rect(slide, 0, 7.15, 13.33, 0.35, NAVY)
    add_text(slide, text, 0.3, 7.17, 12.7, 0.3, font_size=9,
             color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
s1 = blank_slide(prs)
fill_bg(s1, NAVY)

# Gradient-style accent bar left
add_rect(s1, 0, 0, 0.5, 7.5, TEAL)

# Big title
add_text(s1, "Finnie", 1.0, 1.2, 11, 1.6, font_size=80, bold=True,
         color=WHITE, align=PP_ALIGN.LEFT)
add_text(s1, "AI Finance Assistant", 1.0, 2.8, 11, 0.9, font_size=36,
         color=TEAL, bold=False, align=PP_ALIGN.LEFT)

# Divider line
add_rect(s1, 1.0, 3.75, 8.0, 0.05, TEAL)

add_text(s1, "Multi-Agent Conversational System for Financial Education",
         1.0, 4.0, 11, 0.6, font_size=20, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)
add_text(s1, "Powered by Gemini 2.0 Flash · LangGraph · FAISS · Streamlit",
         1.0, 4.7, 11, 0.5, font_size=14, color=DARK_GRAY, align=PP_ALIGN.LEFT)

add_text(s1, "Team Demo  |  April 2026", 1.0, 6.4, 11, 0.5,
         font_size=13, color=DARK_GRAY, italic=True)

footer(s1)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════════════════════════
s2 = blank_slide(prs)
fill_bg(s2, NAVY)
slide_header(s2, "The Problem", "Financial literacy is broken — and expensive to fix")

# Three problem cards
cards = [
    ("💸", "Knowledge Gap",
     "57% of Americans are financially illiterate.\nBasic concepts like compound\ninterest remain unfamiliar."),
    ("🤖", "One-Size-Fits-All",
     "Generic articles and videos can't\nadapt to a beginner's confusion\nor an expert's nuanced question."),
    ("💼", "Advisors Are Inaccessible",
     "Human financial advisors cost\n$150–$400/hr — unreachable\nfor most people who need them most."),
]

for i, (icon, title, body) in enumerate(cards):
    x = 0.6 + i * 4.2
    add_rect(s2, x, 1.55, 3.9, 5.2, SLATE, line_color=TEAL, line_width=1.5)
    add_text(s2, icon, x + 1.5, 1.8, 1.0, 0.9, font_size=36, align=PP_ALIGN.CENTER)
    add_text(s2, title, x + 0.15, 2.8, 3.6, 0.55, font_size=18, bold=True,
             color=TEAL, align=PP_ALIGN.CENTER)
    add_text(s2, body, x + 0.15, 3.45, 3.6, 2.8, font_size=13.5,
             color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

footer(s2)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SOLUTION OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════════
s3 = blank_slide(prs)
fill_bg(s3, NAVY)
slide_header(s3, "The Solution — Finnie", "Personalized financial education through intelligent multi-agent conversation")

add_text(s3, "What Finnie IS", 0.5, 1.5, 6.0, 0.45, font_size=16, bold=True, color=GREEN)
bullet_box(s3, [
    "An adaptive educational AI that explains financial concepts",
    "A multi-agent system (6 specialists) routing each query to the right expert",
    "Context-aware: adjusts language for Beginner / Intermediate / Advanced users",
    "RAG-powered: answers grounded in 35+ curated financial education articles",
    "Free-tier friendly: Gemini 2.0 Flash (60 req/min), yFinance (no key required)",
], 0.5, 2.0, 6.1, 4.5, font_size=14, color=LIGHT_GRAY)

add_text(s3, "What Finnie is NOT", 6.8, 1.5, 6.0, 0.45, font_size=16, bold=True, color=ORANGE)
bullet_box(s3, [
    "Not a financial advisor — never recommends specific investments",
    "Not a trading platform — no buy/sell signals or portfolio management",
    "Not a replacement for CPAs, CFPs, or legal counsel",
    "Not trained on private user financial data",
    "Every response carries context-appropriate educational disclaimers",
], 6.8, 2.0, 6.0, 4.5, font_size=14, color=LIGHT_GRAY)

# Divider
add_rect(s3, 6.55, 1.5, 0.05, 5.3, TEAL)

footer(s3)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════════
s4 = blank_slide(prs)
fill_bg(s4, NAVY)
slide_header(s4, "System Architecture", "Five layers working together — from UI to data")

# Layer boxes
layers = [
    ("STREAMLIT UI  (5 Tabs)",
     "Chat · Portfolio · Market · Goals · Education",
     0.4, SLATE),
    ("LANGGRAPH ORCHESTRATION",
     "profile_loader → intent_classifier → router → agent nodes → synthesizer → disclaimer_injector",
     1.6, SLATE),
    ("6 SPECIALIZED AGENTS",
     "Finance Q&A · Portfolio · Market · Goal Planning · News · Tax",
     2.8, SLATE),
    ("CORE SERVICES",
     "Gemini 2.0 Flash · Rate Limiter (60 req/min) · TTL Cache · Exponential Backoff",
     4.0, SLATE),
    ("DATA LAYER",
     "FAISS IndexFlatIP · 35 KB Articles · Alpha Vantage + yFinance · Google Embeddings",
     5.2, SLATE),
]

for title, body, top, color in layers:
    add_rect(s4, 0.5, top + 1.0, 12.3, 1.0, color, line_color=TEAL, line_width=0.8)
    add_text(s4, title, 0.7, top + 1.1, 4.5, 0.4, font_size=12, bold=True, color=TEAL)
    add_text(s4, body, 5.0, top + 1.1, 7.6, 0.75, font_size=12, color=LIGHT_GRAY)
    if top < 5.2:
        add_text(s4, "▼", 6.3, top + 2.02, 1.0, 0.3, font_size=12, color=TEAL, align=PP_ALIGN.CENTER)

# Arrow labels
add_text(s4, "User Query", 0.55, 1.32, 2.0, 0.25, font_size=10, color=TEAL, italic=True)
add_text(s4, "AgentState (shared)", 0.55, 2.52, 3.0, 0.25, font_size=10, color=TEAL, italic=True)
add_text(s4, "RAG + LLM calls", 0.55, 3.72, 2.5, 0.25, font_size=10, color=TEAL, italic=True)
add_text(s4, "Cached API calls", 0.55, 4.92, 2.5, 0.25, font_size=10, color=TEAL, italic=True)

footer(s4)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — THE SIX AGENTS
# ═══════════════════════════════════════════════════════════════════════════════
s5 = blank_slide(prs)
fill_bg(s5, NAVY)
slide_header(s5, "The Six Specialized Agents", "Each agent is an expert in its domain — routed by intent classification")

agents = [
    ("💬", "Finance Q&A",         "General concepts, definitions,\nhow-X-works questions",       "All KB categories"),
    ("📊", "Portfolio Analysis",  "Allocation %, diversification score,\nrisk vs. tolerance",     "investing, market_concepts"),
    ("📈", "Market Analysis",     "Live index prices, sector data,\nmarket overview",             "market_concepts"),
    ("🎯", "Goal Planning",       "Retirement projections, savings targets,\ncompound interest",  "basics, investing, retirement"),
    ("📰", "News Synthesizer",    "Recent financial news summaries,\neducational context",        "market_concepts"),
    ("🧾", "Tax Education",       "Tax brackets, capital gains,\n401k / IRA / HSA concepts",     "taxes, retirement"),
]

for i, (icon, name, desc, rag) in enumerate(agents):
    row, col = divmod(i, 3)
    x = 0.45 + col * 4.27
    y = 1.55 + row * 2.65
    add_rect(s5, x, y, 3.9, 2.4, SLATE, line_color=TEAL, line_width=1.0)
    add_text(s5, f"{icon}  {name}", x + 0.15, y + 0.1, 3.6, 0.55,
             font_size=15, bold=True, color=TEAL)
    add_text(s5, desc, x + 0.15, y + 0.65, 3.6, 0.9, font_size=13, color=WHITE)
    add_text(s5, f"RAG: {rag}", x + 0.15, y + 1.7, 3.6, 0.45,
             font_size=10.5, color=DARK_GRAY, italic=True)

footer(s5)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — LANGGRAPH WORKFLOW
# ═══════════════════════════════════════════════════════════════════════════════
s6 = blank_slide(prs)
fill_bg(s6, NAVY)
slide_header(s6, "LangGraph Orchestration", "StateGraph with conditional routing and multi-agent fan-out")

# Flow nodes
nodes = [
    ("User Query", 0.3, 3.3, 1.6, 0.6),
    ("profile_loader", 2.1, 3.3, 2.0, 0.6),
    ("intent_classifier", 4.3, 3.3, 2.3, 0.6),
    ("router", 6.8, 3.3, 1.5, 0.6),
]
for label, x, y, w, h in nodes:
    add_rect(s6, x, y, w, h, SLATE, line_color=TEAL, line_width=1.2)
    add_text(s6, label, x, y + 0.12, w, 0.4, font_size=11.5, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s6, "→", x + w, y + 0.15, 0.3, 0.35, font_size=14, color=TEAL, align=PP_ALIGN.CENTER)

# Fan-out agents
agent_nodes = [
    "Finance Q&A", "Portfolio", "Market",
    "Goal Plan", "News", "Tax",
]
for i, name in enumerate(agent_nodes):
    ax = 8.65
    ay = 1.25 + i * 0.87
    add_rect(s6, ax, ay, 2.1, 0.65, SLATE, line_color=GREEN, line_width=1.0)
    add_text(s6, name, ax, ay + 0.1, 2.1, 0.45, font_size=12, color=WHITE,
             align=PP_ALIGN.CENTER, bold=True)

# Bracket line
add_text(s6, "conditional\nfan-out", 8.25, 2.4, 1.2, 0.6, font_size=10, color=TEAL, italic=True)

# Post-synthesis flow
add_text(s6, "→", 10.78, 3.3, 0.4, 0.5, font_size=14, color=TEAL, align=PP_ALIGN.CENTER)
add_rect(s6, 11.1, 3.25, 2.0, 0.7, SLATE, line_color=TEAL, line_width=1.2)
add_text(s6, "response_\nsynthesizer", 11.1, 3.3, 2.0, 0.7, font_size=11,
         bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Bottom row
synthesis_row = [
    ("disclaimer_\ninjector", 3.5, 5.5),
    ("profile_\nupdater", 6.2, 5.5),
    ("END", 8.8, 5.5),
]
prev_x = 11.1 + 1.0
add_text(s6, "↓", 12.0, 4.0, 0.5, 0.5, font_size=16, color=TEAL, align=PP_ALIGN.CENTER)

for label, x, y in synthesis_row:
    add_rect(s6, x, y, 2.2, 0.7, SLATE, line_color=TEAL, line_width=1.2)
    add_text(s6, label, x, y + 0.05, 2.2, 0.65, font_size=11.5,
             bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s6, "←───────────────────────────────────────────────────────────────", 3.0, 5.1, 10.0, 0.4,
         font_size=11, color=TEAL)

# AgentState callout
add_rect(s6, 0.3, 1.2, 1.7, 1.8, NAVY, line_color=ORANGE, line_width=1.0)
add_text(s6, "AgentState\n(shared)\n\n• query\n• intent\n• outputs\n• profile",
         0.35, 1.25, 1.6, 1.7, font_size=9.5, color=LIGHT_GRAY)

footer(s6)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — RAG PIPELINE
# ═══════════════════════════════════════════════════════════════════════════════
s7 = blank_slide(prs)
fill_bg(s7, NAVY)
slide_header(s7, "RAG Pipeline", "Retrieval-Augmented Generation — grounding answers in curated knowledge")

# Left: Knowledge base
add_rect(s7, 0.4, 1.45, 3.5, 5.6, SLATE, line_color=TEAL, line_width=1.0)
add_text(s7, "📚  Knowledge Base", 0.55, 1.55, 3.2, 0.5, font_size=14, bold=True, color=TEAL)
kb_items = [
    "basics/          (8 articles)",
    "investing/       (9 articles)",
    "retirement/      (5 articles)",
    "taxes/           (5 articles)",
    "budgeting/       (3 articles)",
    "market_concepts/ (5 articles)",
    "",
    "35 total  ·  YAML frontmatter",
    "512-token chunks  ·  50 overlap",
]
bullet_box(s7, kb_items, 0.55, 2.1, 3.2, 4.5, font_size=12.5, bullet="•", color=LIGHT_GRAY, spacing=20)

# Middle: Pipeline steps
steps = [
    ("1", "ArticleChunker", "512-token chunks with\n50-token overlap + metadata"),
    ("2", "FAISSIndexer", "IndexFlatIP + L2-normalize\n→ cosine similarity"),
    ("3", "EmbeddingModel", "Google text-embedding-004\n(HuggingFace MiniLM fallback)"),
    ("4", "FAISSRetriever", "Query embed → top-K search\n→ MMR re-ranking → filter"),
    ("5", "RAGPipeline", "Format context string\n+ source citations + TTL cache"),
]
for i, (num, name, desc) in enumerate(steps):
    y = 1.45 + i * 1.1
    add_rect(s7, 4.2, y, 0.5, 0.85, TEAL)
    add_text(s7, num, 4.2, y + 0.15, 0.5, 0.55, font_size=16, bold=True,
             color=NAVY, align=PP_ALIGN.CENTER)
    add_rect(s7, 4.85, y, 3.8, 0.85, SLATE, line_color=TEAL, line_width=0.8)
    add_text(s7, name, 5.0, y + 0.03, 3.5, 0.4, font_size=13, bold=True, color=WHITE)
    add_text(s7, desc, 5.0, y + 0.43, 3.5, 0.5, font_size=10.5, color=DARK_GRAY)

# Right: Output
add_rect(s7, 9.0, 1.45, 3.9, 5.6, SLATE, line_color=GREEN, line_width=1.0)
add_text(s7, "✅  Agent Receives", 9.15, 1.55, 3.6, 0.5, font_size=14, bold=True, color=GREEN)
output_items = [
    "Formatted context string",
    "[Source N: title (category)]",
    "Top-5 relevant chunks",
    "Deduplicated citations",
    "Cached for 1-hour TTL",
    "",
    "Agent injects into",
    "LLM system prompt →",
    "Grounded, cited response",
]
bullet_box(s7, output_items, 9.15, 2.15, 3.6, 4.5, font_size=12.5, bullet="✓", color=LIGHT_GRAY, spacing=20)

footer(s7)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — TECH STACK
# ═══════════════════════════════════════════════════════════════════════════════
s8 = blank_slide(prs)
fill_bg(s8, NAVY)
slide_header(s8, "Tech Stack", "Production-grade components, all open-source or free-tier")

stack = [
    ("🤖", "LLM",             "Google Gemini 2.0 Flash",             "60 req/min free tier · temperature tuned per use-case"),
    ("🔗", "Orchestration",   "LangGraph StateGraph",                 "Conditional edges · multi-agent fan-out · shared state"),
    ("🗄️", "Vector DB",       "FAISS IndexFlatIP",                    "Local · cosine similarity · MMR re-ranking"),
    ("🔢", "Embeddings",      "Google text-embedding-004",            "HuggingFace MiniLM fallback (no API key required)"),
    ("📊", "Market Data",     "Alpha Vantage + yFinance",             "25 req/day (AV) · unlimited fallback (yF) · 30-min cache"),
    ("🖥️", "UI",              "Streamlit + Plotly",                   "5 tabs · streaming chat · CSV upload · interactive charts"),
    ("⚡", "Rate Limiting",   "Token Bucket (60 req/min)",            "Thread-safe · asyncio-compatible · exponential backoff"),
    ("💾", "Caching",         "In-memory TTL + LRU",                  "Market: 30 min · News: 15 min · LLM: 24 hr · RAG: 1 hr"),
    ("🧪", "Testing",         "pytest + pytest-cov",                  "166 tests · unit + integration · ≥70% coverage enforced"),
]

for i, (icon, layer, tech, detail) in enumerate(stack):
    y = 1.4 + i * 0.65
    add_rect(s8, 0.4, y, 0.8, 0.55, SLATE)
    add_text(s8, icon, 0.4, y + 0.05, 0.8, 0.45, font_size=18, align=PP_ALIGN.CENTER)
    add_text(s8, layer, 1.3, y + 0.07, 1.8, 0.45, font_size=12, bold=True, color=TEAL)
    add_text(s8, tech,  3.2, y + 0.07, 3.5, 0.45, font_size=12, bold=True, color=WHITE)
    add_text(s8, detail, 6.8, y + 0.07, 6.2, 0.45, font_size=11, color=DARK_GRAY)

footer(s8)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — REGULATORY / DISCLAIMER DESIGN
# ═══════════════════════════════════════════════════════════════════════════════
s9 = blank_slide(prs)
fill_bg(s9, NAVY)
slide_header(s9, "Responsible Design", "Education vs. advice — built-in by architecture, not afterthought")

# Left column: What we do
add_rect(s9, 0.4, 1.45, 5.9, 5.6, SLATE, line_color=GREEN, line_width=1.2)
add_text(s9, "✅  What Finnie Does", 0.6, 1.55, 5.5, 0.5, font_size=15, bold=True, color=GREEN)
dos = [
    "Explains financial concepts clearly",
    "Adapts language to user knowledge level",
    "Cites every answer with source articles",
    "Injects context-appropriate disclaimers",
    "Labels all projections as illustrative",
    "Timestamps all market data",
    "Recommends professional consultation",
    "Logs every response for auditability",
]
bullet_box(s9, dos, 0.6, 2.1, 5.6, 4.7, font_size=13.5, bullet="✓", color=LIGHT_GRAY, spacing=21)

# Right column: What we don't
add_rect(s9, 6.65, 1.45, 6.2, 5.6, SLATE, line_color=ORANGE, line_width=1.2)
add_text(s9, "🚫  What Finnie Never Does", 6.85, 1.55, 5.9, 0.5, font_size=15, bold=True, color=ORANGE)
donts = [
    "Recommend specific stocks or funds",
    "Provide buy/sell signals",
    "Guarantee investment returns",
    "File or review tax returns",
    "Manage a user's portfolio",
    "Store personal financial data",
    "Act as a licensed financial advisor",
    "Make predictions about market direction",
]
bullet_box(s9, donts, 6.85, 2.1, 5.9, 4.7, font_size=13.5, bullet="✗", color=LIGHT_GRAY, spacing=21)

# Disclaimer box
add_rect(s9, 0.4, 7.0, 12.5, 0.0, NAVY)  # spacer
footer(s9, "Disclaimer system: 5 context-aware disclaimers injected by dedicated graph node — general · portfolio · tax · market · projections")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — LIVE DEMO GUIDE
# ═══════════════════════════════════════════════════════════════════════════════
s10 = blank_slide(prs)
fill_bg(s10, NAVY)
slide_header(s10, "Live Demo", "6 segments · ~8 minutes · streamlit run run.py")

segments = [
    ("1", "Onboarding",           "1 min",  "Select knowledge level (Beginner), risk tolerance (Conservative), goal (Retirement 30 yrs)"),
    ("2", "Multi-turn Chat",       "2 min",  '"What is an index fund?" → follow-up → "How does it differ from active funds?" — show context memory'),
    ("3", "Portfolio Analysis",    "2 min",  "Upload CSV (AAPL/MSFT/BND/VTI) → pie chart → allocation analysis → no buy/sell advice"),
    ("4", "Market Data",           "1 min",  'Market tab → live S&P 500 / NASDAQ → return to chat: "How did tech perform today?"'),
    ("5", "Goal Planning",         "1 min",  '$1M retirement, 30 yr, $50k savings, $1k/month → 3 projection scenarios (5%/7%/9%) → "educational illustration"'),
    ("6", "Multi-Agent Query",     "1 min",  '"How do my tech gains affect my taxes?" → Portfolio + Tax agents both fire → merged response'),
]

for i, (num, title, dur, script) in enumerate(segments):
    y = 1.45 + i * 0.97
    add_rect(s10, 0.35, y, 0.65, 0.82, TEAL)
    add_text(s10, num, 0.35, y + 0.15, 0.65, 0.5, font_size=20, bold=True,
             color=NAVY, align=PP_ALIGN.CENTER)
    add_rect(s10, 1.1, y, 3.0, 0.82, SLATE, line_color=TEAL, line_width=0.8)
    add_text(s10, title, 1.2, y + 0.03, 2.0, 0.42, font_size=13, bold=True, color=WHITE)
    add_text(s10, dur, 2.9, y + 0.03, 1.0, 0.42, font_size=12, color=TEAL, align=PP_ALIGN.RIGHT)
    add_text(s10, script, 4.25, y + 0.1, 8.8, 0.72, font_size=11.5, color=LIGHT_GRAY)

footer(s10)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — MILESTONE ROADMAP
# ═══════════════════════════════════════════════════════════════════════════════
s11 = blank_slide(prs)
fill_bg(s11, NAVY)
slide_header(s11, "Milestone Roadmap", "5 milestones · 30-day delivery · all completed ✅")

milestones = [
    ("M1", "Foundation",          "Days 1–5",   "Config · LLM client · rate limiter · cache · retry · base agent",              True),
    ("M2", "RAG Pipeline",        "Days 5–10",  "35 articles · FAISS index · chunker · retriever · MMR pipeline",               True),
    ("M3", "All 6 Agents + Graph","Days 10–18", "6 agents · LangGraph StateGraph · intent classifier · router · nodes",         True),
    ("M4", "Streamlit UI",        "Days 18–24", "5 tabs · streaming chat · CSV upload · Plotly charts · onboarding",            True),
    ("M5", "Tests + Docs",        "Days 24–30", "166 tests · ≥70% coverage · README · bug fixes · production polish",           True),
]

# Timeline bar
add_rect(s11, 0.5, 2.1, 12.3, 0.12, TEAL)
for i in range(5):
    x = 0.5 + i * 2.46
    add_rect(s11, x, 1.85, 0.12, 0.6, TEAL)

for i, (code, title, days, detail, done) in enumerate(milestones):
    x = 0.45 + i * 2.46
    y_top = 1.25

    # Milestone marker
    color = GREEN if done else ORANGE
    add_rect(s11, x - 0.1, 1.78, 0.7, 0.7, color)
    add_text(s11, "✅" if done else "🔄", x - 0.1, 1.82, 0.7, 0.5,
             font_size=16, align=PP_ALIGN.CENTER)

    # Card below
    add_rect(s11, x - 0.25, 2.35, 2.35, 4.6, SLATE, line_color=color, line_width=1.0)
    add_text(s11, code, x - 0.1, 2.42, 2.0, 0.4, font_size=14, bold=True, color=color)
    add_text(s11, title, x - 0.1, 2.82, 2.0, 0.5, font_size=12.5, bold=True, color=WHITE)
    add_text(s11, days, x - 0.1, 3.3, 2.0, 0.35, font_size=11, color=TEAL, italic=True)
    add_text(s11, detail, x - 0.1, 3.72, 2.15, 3.0, font_size=10, color=DARK_GRAY)

footer(s11)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Q&A / CLOSING
# ═══════════════════════════════════════════════════════════════════════════════
s12 = blank_slide(prs)
fill_bg(s12, NAVY)

add_rect(s12, 0, 0, 0.5, 7.5, TEAL)

add_text(s12, "Questions?", 1.0, 1.5, 11, 1.4, font_size=70, bold=True,
         color=WHITE, align=PP_ALIGN.LEFT)
add_rect(s12, 1.0, 3.1, 8.0, 0.06, TEAL)

add_text(s12, "Try it yourself:", 1.0, 3.4, 11, 0.5, font_size=18, bold=True, color=TEAL)
add_text(s12, "streamlit run run.py", 1.0, 3.95, 9, 0.55,
         font_size=22, bold=True, color=WHITE)

key_points = [
    "166 tests passing  ·  ≥70% coverage enforced",
    "6 specialized agents  ·  LangGraph StateGraph",
    "35 curated KB articles  ·  FAISS MMR retrieval",
    "Educational only — never financial advice",
]
bullet_box(s12, key_points, 1.0, 5.1, 10, 1.8, font_size=15, bullet="▸",
           color=LIGHT_GRAY, spacing=24)

footer(s12, "Finnie — AI Finance Assistant  |  Educational AI, not financial advice  |  Team Demo  April 2026")


# ── Save ──────────────────────────────────────────────────────────────────────
out_path = "Finnie_AI_Finance_Assistant_Demo.pptx"
prs.save(out_path)
print(f"✅  Saved: {out_path}")
print(f"   Slides: {len(prs.slides)}")
